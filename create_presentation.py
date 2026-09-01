import os
import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def update_presentation():
    prs = pptx.Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Harmonious Course-Themed Palette
    DARK_BLUE = RGBColor(16, 37, 66)      # Primary Dark #102542
    ACCENT_TEAL = RGBColor(23, 162, 184)  # Teal Accent
    ACCENT_BLUE = RGBColor(41, 98, 153)   # Accent Blue
    LIGHT_BG = RGBColor(248, 249, 250)    # Soft Gray Card
    TEXT_DARK = RGBColor(33, 37, 41)      # Body text
    TEXT_MUTED = RGBColor(108, 117, 125)  # Muted subtitles
    WHITE = RGBColor(255, 255, 255)
    BORDER_COLOR = RGBColor(220, 225, 230)
    GREEN_CARD = RGBColor(235, 247, 238)
    GREEN_TEXT = RGBColor(40, 167, 69)

    def add_header(slide, title_text, phase_text="DATA ANALYTICS PROJECT WORK"):
        header_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2))
        header_box.fill.solid()
        header_box.fill.fore_color.rgb = DARK_BLUE
        header_box.line.color.rgb = DARK_BLUE
        
        tf_cat = header_box.text_frame
        tf_cat.word_wrap = True
        tf_cat.margin_left = Inches(0.8)
        tf_cat.margin_top = Inches(0.15)
        
        p0 = tf_cat.paragraphs[0]
        p0.text = phase_text.upper()
        p0.font.size = Pt(11)
        p0.font.bold = True
        p0.font.color.rgb = RGBColor(180, 215, 235)
        
        p1 = tf_cat.add_paragraph()
        p1.text = title_text
        p1.font.size = Pt(22)
        p1.font.bold = True
        p1.font.color.rgb = WHITE

    # --- SLIDE 1: Title Slide ---
    slide1 = prs.slides.add_slide(blank_layout)
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = DARK_BLUE
    bg1.line.color.rgb = DARK_BLUE

    tbox1 = slide1.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.333), Inches(4.0))
    tf1 = tbox1.text_frame
    tf1.word_wrap = True
    
    p = tf1.paragraphs[0]
    p.text = "Data Analytics & Predictive Modeling"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    p2 = tf1.add_paragraph()
    p2.text = "Flue Gas Emissions (CO, NOx) and Energy Yield in Gas Turbines"
    p2.font.size = Pt(22)
    p2.font.color.rgb = RGBColor(180, 215, 235)
    
    p3 = tf1.add_paragraph()
    p3.text = "\nCourse: Data Analytics (and Data Driven Decision) — University of L'Aquila\nBased on Multi-Year Hourly Telemetry (36,733 records) and Course Methodologies"
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(200, 215, 230)

    # --- SLIDE 2: Project Workflow & 7 Steps ---
    slide2 = prs.slides.add_slide(blank_layout)
    add_header(slide2, "Project Workflow & Guidelines Structure", "METHODOLOGICAL ROADMAP")
    
    steps = [
        ("1. Description of Dataset", "Data origin, 11 sensor variables, 36,733 hourly observations, train/test split protocol."),
        ("2. Data Cleaning & Validation", "Verification of missing values, range validation, physical plausibility check, standard scaling."),
        ("3. Exploratory Data Analysis", "Descriptive statistics, distribution shapes, correlation heatmaps, environmental impact."),
        ("4. Main Analysis & Methods", "Techniques aligned with course lectures: Dimensionality Reduction, Regression, Diagnostics."),
        ("5. Preview of Results", "High-level summary of model performance, key feature importances, emission trade-offs."),
        ("6. Detailed Results", "In-depth metric breakdown (R², RMSE, residuals), cross-validation across 2011–2015."),
        ("7. Conclusions", "Industrial takeaways, operational recommendations, turbine emission control strategies.")
    ]
    for i, (title, desc) in enumerate(steps):
        col = 0 if i < 4 else 1
        row = i if i < 4 else i - 4
        left = Inches(0.8 + col * 6.0)
        top = Inches(1.5 + row * 1.35)
        card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.7), Inches(1.2))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_BG
        card.line.color.rgb = BORDER_COLOR
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.15)
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE
        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(11)
        p_desc.font.color.rgb = TEXT_DARK

    # --- SLIDE 3: 1. Description of the Dataset ---
    slide3 = prs.slides.add_slide(blank_layout)
    add_header(slide3, "1. Description of the Dataset & Sensor Features", "PHASE 1: DATASET PROFILING")
    
    card_left = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(4.0), Inches(5.5))
    card_left.fill.solid()
    card_left.fill.fore_color.rgb = LIGHT_BG
    card_left.line.color.rgb = BORDER_COLOR
    tf_l = card_left.text_frame
    tf_l.word_wrap = True
    tf_l.margin_left = Inches(0.25)
    tf_l.margin_top = Inches(0.25)
    p = tf_l.paragraphs[0]
    p.text = "Dataset Characteristics"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    info_points = [
        ("Location:", "Gas turbine plant in NW Turkey"),
        ("Timeframe:", "2011 – 2015 (5 consecutive years)"),
        ("Granularity:", "Hourly aggregated averages/sums"),
        ("Observations:", "36,733 instances (sorted in time)"),
        ("Missing Values:", "0 (Complete sensor records)"),
        ("Data Protocol:", "Train: 2011-13 (22,191 samples, 60.4%)\nTest: 2014-15 (14,535 samples, 39.6%)"),
        ("Key Targets:", "Emissions (CO, NOx), Energy (TEY)")
    ]
    for lbl, val in info_points:
        p_item = tf_l.add_paragraph()
        p_item.text = f"• {lbl} {val}"
        p_item.font.size = Pt(12)
        p_item.font.color.rgb = TEXT_DARK
        
    table_shape = slide3.shapes.add_table(12, 4, Inches(5.1), Inches(1.5), Inches(7.4), Inches(5.5))
    table = table_shape.table
    table.columns[0].width = Inches(1.5)
    table.columns[1].width = Inches(1.1)
    table.columns[2].width = Inches(1.0)
    table.columns[3].width = Inches(3.8)
    headers = ["Category", "Variable", "Unit", "Physical Role / Meaning"]
    for col_idx, h in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = WHITE
    sensor_data = [
        ("Ambient", "AT", "°C", "Ambient Temperature (inlet air density driver)"),
        ("Ambient", "AP", "mbar", "Ambient Atmospheric Pressure"),
        ("Ambient", "AH", "%", "Ambient Relative Humidity"),
        ("Turbine", "AFDP", "mbar", "Air Filter Difference Pressure (filter clogging)"),
        ("Turbine", "GTEP", "mbar", "Gas Turbine Exhaust Pressure"),
        ("Turbine", "TIT", "°C", "Turbine Inlet Temperature (efficiency driver)"),
        ("Turbine", "TAT", "°C", "Turbine After Temperature (exhaust)"),
        ("Turbine", "CDP", "mbar", "Compressor Discharge Pressure"),
        ("Yield", "TEY", "MWh", "Turbine Energy Yield (Net power output)"),
        ("Emission", "CO", "mg/m³", "Carbon Monoxide (Incomplete combustion)"),
        ("Emission", "NOx", "mg/m³", "Nitrogen Oxides (Thermal NOx mechanism)")
    ]
    for row_idx, row in enumerate(sensor_data, start=1):
        for col_idx, val in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if row_idx % 2 == 0 else LIGHT_BG
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(10)
            p.font.color.rgb = TEXT_DARK

    # --- SLIDE 4: 2. Data Cleaning & Preprocessing ---
    slide4 = prs.slides.add_slide(blank_layout)
    add_header(slide4, "2. Data Cleaning, Quality Validation & Preprocessing", "PHASE 1: DATA PREPARATION")
    card_c1 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(3.7), Inches(2.5))
    card_c1.fill.solid()
    card_c1.fill.fore_color.rgb = LIGHT_BG
    card_c1.line.color.rgb = BORDER_COLOR
    tf_c1 = card_c1.text_frame
    tf_c1.word_wrap = True
    tf_c1.margin_left = Inches(0.2)
    tf_c1.margin_top = Inches(0.2)
    p = tf_c1.paragraphs[0]
    p.text = "1. Completeness & Integrity"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    for it in ["• Missing Values: 0 nulls across all 11 columns.", "• Duplicates: 7 sensor-profile duplicates pruned (36,726 clean rows).", "• Chronology: Preserved time order across 2011–2015.", "• Sanity Checks: Range verified for positive pressures & valid temps."]:
        p = tf_c1.add_paragraph()
        p.text = it
        p.font.size = Pt(10.5)
        p.font.color.rgb = TEXT_DARK

    card_c2 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.8), Inches(1.5), Inches(3.7), Inches(2.5))
    card_c2.fill.solid()
    card_c2.fill.fore_color.rgb = LIGHT_BG
    card_c2.line.color.rgb = BORDER_COLOR
    tf_c2 = card_c2.text_frame
    tf_c2.word_wrap = True
    tf_c2.margin_left = Inches(0.2)
    tf_c2.margin_top = Inches(0.2)
    p = tf_c2.paragraphs[0]
    p.text = "2. σ-Clipping (Course Method)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    for it in ["• Robust Dispersion: $\\hat{\\sigma} = \\frac{IQR}{1.35} = \\frac{Q_3 - Q_1}{1.35}$", "• Robust Centering: Uses Median ($Q_2$) to resist leverage.", "• Boundaries: $[\\text{Med} - 5\\hat{\\sigma},\\; \\text{Med} + 5\\hat{\\sigma}]$", "• Ambient/Turbine: 0 outliers (100% physically valid).", "• Inspection Only: Extreme values retained as real transient conditions."]:
        p = tf_c2.add_paragraph()
        p.text = it
        p.font.size = Pt(10.5)
        p.font.color.rgb = TEXT_DARK

    card_c3 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), Inches(1.5), Inches(3.7), Inches(2.5))
    card_c3.fill.solid()
    card_c3.fill.fore_color.rgb = LIGHT_BG
    card_c3.line.color.rgb = BORDER_COLOR
    tf_c3 = card_c3.text_frame
    tf_c3.word_wrap = True
    tf_c3.margin_left = Inches(0.2)
    tf_c3.margin_top = Inches(0.2)
    p = tf_c3.paragraphs[0]
    p.text = "3. Protocol & Scaling"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    for it in ["• Train Set (2011–13): 22,191 samples (60.4%)", "• Test Set (2014–15): 14,535 samples (39.6%)", "• No Leakage: Fit scaler strictly on Train, transform Test.", "• Target Separation: Standardized 8 predictors and 3 targets independently."]:
        p = tf_c3.add_paragraph()
        p.text = it
        p.font.size = Pt(10.5)
        p.font.color.rgb = TEXT_DARK

    table_shape4 = slide4.shapes.add_table(7, 6, Inches(0.8), Inches(4.3), Inches(11.7), Inches(2.7))
    t4 = table_shape4.table
    for c in range(6):
        t4.columns[c].width = Inches(1.95)
    headers4 = ["Variable", "Median (Q2)", "IQR (Q3 - Q1)", "Est. σ̂ (IQR/1.35)", "Valid Interval [Med ± 5σ̂]", "Physical Assessment"]
    for col_idx, h in enumerate(headers4):
        cell = t4.cell(0, col_idx)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(10)
        p.font.color.rgb = WHITE
    sigma_data = [
        ("AT (°C)", "17.80", "11.88", "8.80", "[-26.22, 61.82]", "0 outliers (Physical weather)"),
        ("AP (mbar)", "1012.60", "8.20", "6.07", "[982.23, 1042.97]", "0 outliers (Barometric band)"),
        ("TIT (°C)", "1085.90", "25.20", "18.67", "[992.57, 1179.23]", "0 outliers (Firing regime)"),
        ("TEY (MWh)", "133.73", "19.63", "14.54", "[61.03, 206.43]", "0 outliers (Generator range)"),
        ("CO (mg/m³)", "1.71", "1.66", "1.23", "[-4.44, 7.87]", "Spikes = Incomplete combustion"),
        ("NOx (mg/m³)", "63.85", "14.38", "10.66", "[10.57, 117.13]", "Peaks = High flame temp")
    ]
    for row_idx, row in enumerate(sigma_data, start=1):
        for col_idx, val in enumerate(row):
            cell = t4.cell(row_idx, col_idx)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if row_idx % 2 == 0 else LIGHT_BG
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(9.5)
            p.font.color.rgb = TEXT_DARK

    # --- SLIDE 5: 3. EDA - Distributions & Correlation ---
    slide5 = prs.slides.add_slide(blank_layout)
    add_header(slide5, "3. Exploratory Analysis: Distributions & Correlation Structure", "PHASE 2: EXPLORATORY DATA ANALYSIS")
    card_eda1 = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.4), Inches(5.6), Inches(5.6))
    card_eda1.fill.solid()
    card_eda1.fill.fore_color.rgb = LIGHT_BG
    card_eda1.line.color.rgb = BORDER_COLOR
    tf_e1 = card_eda1.text_frame
    tf_e1.word_wrap = True
    tf_e1.margin_left = Inches(0.25)
    tf_e1.margin_top = Inches(0.2)
    p = tf_e1.paragraphs[0]
    p.text = "Course Definition & Statistical Synthesis"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    eda_points = [
        ("Exploratory Analysis (Slide 3 & 10):", "Unsupervised synthesis of distributions into Size (Mean $\\mu$, Median, Mode) and Spread (Variance $\\sigma^2$, Std, IQR, Gini $G$, Entropy $H$)."),
        ("Size vs Spread Synthesis:", "• Mean vs Median: CO mean (2.21) >> median (1.52) demonstrates strong positive skew.\n• Gini Index: CO exhibits high concentration (G=0.442) due to episodic spikes."),
        ("Key Pearson Correlation ($r_{XY}$) Findings:", "• Compressor & Energy Coupling: CDP & TEY ($r = 0.99$), GTEP & TEY ($r = 0.98$) indicate strong collinearity in operating load.\n• Environmental Cooling Effect: Ambient Temp AT vs TEY ($r = -0.58$) due to reduced air density at higher temperatures.\n• Combustion Trade-off: CO vs NOx ($r = -0.37$) reflects the physical inverse relationship between complete oxidation and thermal NOx.")
    ]
    for lbl, desc in eda_points:
        p_l = tf_e1.add_paragraph()
        p_l.text = f"\n{lbl}"
        p_l.font.bold = True
        p_l.font.size = Pt(11)
        p_l.font.color.rgb = ACCENT_BLUE
        p_d = tf_e1.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(10)
        p_d.font.color.rgb = TEXT_DARK
    if os.path.exists("figures/eda_correlation_heatmap.png"):
        slide5.shapes.add_picture("figures/eda_correlation_heatmap.png", Inches(6.7), Inches(1.4), width=Inches(5.8))

    # --- SLIDE 6: 3. EDA - PCA & Physical Insights ---
    slide6 = prs.slides.add_slide(blank_layout)
    add_header(slide6, "3. Exploratory Analysis: PCA Dimensionality Reduction & Biplot", "PHASE 2: UNSUPERVISED PCA")
    card_pca = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.4), Inches(5.6), Inches(5.6))
    card_pca.fill.solid()
    card_pca.fill.fore_color.rgb = LIGHT_BG
    card_pca.line.color.rgb = BORDER_COLOR
    tf_pca = card_pca.text_frame
    tf_pca.word_wrap = True
    tf_pca.margin_left = Inches(0.25)
    tf_pca.margin_top = Inches(0.2)
    p = tf_pca.paragraphs[0]
    p.text = "PCA as an Exploratory Technique (Slide 70)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    pca_notes = [
        ("Course Formulation (Slide 75):", "Eigen-decomposition of Covariance Matrix $C v_i = \\lambda_i v_i$.\nProportion of Variance Explained: $\\text{PVE}_i = \\lambda_i / \\sum \\lambda_j$."),
        ("Dimensionality Reduction Results:", "• PC1 (48.28%): Captures general turbine thermodynamic load (high loadings for CDP, TEY, GTEP, TIT).\n• PC2 (20.48%): Captures ambient temperature (AT) and the inverse CO vs NOx combustion trade-off.\n• First 2 PCs capture 68.76% of total system variance.\n• First 5 PCs capture 92.02% (Scree plot elbow at $p=5$)."),
        ("Industrial Takeaway:", "The 11 physical sensors effectively lie on a low-dimensional manifold governed by ambient weather and power dispatch.")
    ]
    for lbl, desc in pca_notes:
        p_l = tf_pca.add_paragraph()
        p_l.text = f"\n{lbl}"
        p_l.font.bold = True
        p_l.font.size = Pt(11)
        p_l.font.color.rgb = ACCENT_BLUE
        p_d = tf_pca.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(10)
        p_d.font.color.rgb = TEXT_DARK
    if os.path.exists("figures/eda_pca_biplot.png"):
        slide6.shapes.add_picture("figures/eda_pca_biplot.png", Inches(6.7), Inches(1.4), width=Inches(5.8))

    # --- SLIDE 7: 4. Main Analysis Objectives ---
    slide7 = prs.slides.add_slide(blank_layout)
    add_header(slide7, "4. Main Analysis: Objectives & Methodological Framework", "PHASE 3: METHODOLOGICAL FRAMEWORK")
    objs = [
        ("Objective 1: Energy Yield (TEY) & PCR", 
         "• Predict net turbine energy yield from 8 sensor telemetry features.\n"
         "• Problem: Extreme Multicollinearity (VIF > 250 for CDP, TIT, GTEP, Slide 115).\n"
         "• Method: Principal Component Regression (PCR, Slides 70, 106) on 4 feature PCs ($Z = X_{\\text{std}} V_4$, 90.76% PVE) to eliminate variance inflation and ensure numerical stability."),
        ("Objective 2: Emissions Modeling (CO & NOx)", 
         "• Quantify pollutant generation as a function of firing state.\n"
         "• Method: Polynomial Regression (Slide 116 & Ex 4.1) for CO to capture steep non-linear spikes during low-load/startup regimes ($x \\to [x, x^2]$).\n"
         "• Multivariate OLS for NOx tracking thermal oxidation."),
        ("Objective 3: Operational Regimes (K-Means)", 
         "• Discover discrete turbine operational regimes (Unsupervised).\n"
         "• Method: K-Means Clustering (Slides 90-91) with Silhouette validation (Slide 94, Ex 3.1): $\\text{SIL} = \\frac{1}{m}\\sum \\frac{b_i - a_i}{\\max(a_i, b_i)}$.\n"
         "• Identifies Base-load, Peak-load, and Low-load emission profiles.")
    ]
    for i, (title, desc) in enumerate(objs):
        left = Inches(0.8 + i * 4.0)
        card_obj = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.5), Inches(3.7), Inches(3.2))
        card_obj.fill.solid()
        card_obj.fill.fore_color.rgb = LIGHT_BG
        card_obj.line.color.rgb = BORDER_COLOR
        tf_o = card_obj.text_frame
        tf_o.word_wrap = True
        tf_o.margin_left = Inches(0.2)
        tf_o.margin_top = Inches(0.2)
        p = tf_o.paragraphs[0]
        p.text = title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE
        p_desc = tf_o.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(10)
        p_desc.font.color.rgb = TEXT_DARK

    bot_card = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.9), Inches(11.7), Inches(2.1))
    bot_card.fill.solid()
    bot_card.fill.fore_color.rgb = LIGHT_BG
    bot_card.line.color.rgb = BORDER_COLOR
    tf_b = bot_card.text_frame
    tf_b.word_wrap = True
    tf_b.margin_left = Inches(0.25)
    tf_b.margin_top = Inches(0.15)
    p = tf_b.paragraphs[0]
    p.text = "Mathematical Formulations & Accuracy Metrics (Course Standards)"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p_f = tf_b.add_paragraph()
    p_f.text = (
        "• Ordinary Least Squares: $\\hat{\\beta} = (X^T X)^{-1} X^T y$ (Slide 106, computed via `np.linalg.lstsq` for numerical stability)\n"
        "• Principal Component Regression: $\\hat{\\beta}_{PCR} = (Z^T Z)^{-1} Z^T y$, where $Z = X_{\\text{std}} V_4$ (4 PCs capturing 90.76% PVE)\n"
        "• Multicollinearity VIF: $\\text{VIF}(\\hat{\\beta}_j) = \\frac{1}{1 - R^2_{x_j|x_{-j}}}$ (Slide 115) | Silhouette Quality: $\\text{SIL}_i = \\frac{b_i - a_i}{\\max(a_i, b_i)}$ (Slide 94)"
    )
    p_f.font.size = Pt(10.5)
    p_f.font.color.rgb = TEXT_DARK

    # --- SLIDE 8: 5. Preview & Summary of Results ---
    slide8 = prs.slides.add_slide(blank_layout)
    add_header(slide8, "5. Preview & High-Level Summary of Results", "PHASE 4: RESULTS SYNTHESIS")

    # Top 4 KPI Stat Callout Cards
    kpis = [
        ("0.9501", "Best TEY Test R² (OLS)", "Full 8-feature multivariate model"),
        ("0.9397", "TEY PCR Test R² (4 PCs)", "Collinearity-free (90.76% PVE)"),
        ("0.6477", "CO Train R² (Poly Deg. 2)", "Captures non-linear low-load curve"),
        ("3 Regimes", "Operational Clusters", "Peak, Baseload & Part-Load (K-Means)")
    ]
    for i, (val, lbl, sub) in enumerate(kpis):
        left = Inches(0.8 + i * 3.0)
        card_kpi = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.4), Inches(2.75), Inches(1.6))
        card_kpi.fill.solid()
        card_kpi.fill.fore_color.rgb = GREEN_CARD if i < 2 else LIGHT_BG
        card_kpi.line.color.rgb = GREEN_TEXT if i < 2 else BORDER_COLOR
        
        tf_k = card_kpi.text_frame
        tf_k.word_wrap = True
        tf_k.margin_left = Inches(0.15)
        tf_k.margin_top = Inches(0.12)
        
        p = tf_k.paragraphs[0]
        p.text = val
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = GREEN_TEXT if i < 2 else DARK_BLUE
        
        p_l = tf_k.add_paragraph()
        p_l.text = lbl
        p_l.font.size = Pt(10.5)
        p_l.font.bold = True
        p_l.font.color.rgb = TEXT_DARK
        
        p_s = tf_k.add_paragraph()
        p_s.text = sub
        p_s.font.size = Pt(9)
        p_s.font.color.rgb = TEXT_MUTED

    # Bottom Master Table
    t_shape8 = slide8.shapes.add_table(8, 5, Inches(0.8), Inches(3.2), Inches(11.7), Inches(3.8))
    t8 = t_shape8.table
    t8.columns[0].width = Inches(3.7)
    t8.columns[1].width = Inches(2.0)
    t8.columns[2].width = Inches(2.0)
    t8.columns[3].width = Inches(2.0)
    t8.columns[4].width = Inches(2.0)

    headers8 = ["Target Variable & Model Specification", "Train R² (2011–13)", "Train RMSE", "Test R² (2014–15)", "Test RMSE"]
    for c_idx, h in enumerate(headers8):
        cell = t8.cell(0, c_idx)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = WHITE

    model_summary_data = [
        ("TEY: Full Multivariate OLS (8 Features)", "0.9977", "0.76 MWh", "0.9501", "3.35 MWh"),
        ("TEY: PCR (4 Principal Components, 90.76% PVE)", "0.9693", "2.81 MWh", "0.9397", "3.68 MWh"),
        ("TEY: Training-Mean Baseline", "0.0000", "16.03 MWh", "-0.0000", "14.98 MWh"),
        ("CO: Polynomial (Degree 2, Slide 116)", "0.6477", "1.36 mg/m³", "0.2635", "1.88 mg/m³"),
        ("CO: Training-Mean Baseline", "0.0000", "2.30 mg/m³", "-0.0333", "2.23 mg/m³"),
        ("NOx: Multivariate OLS (8 Features)", "0.4536", "8.16 mg/m³", "-1.0928*", "15.30 mg/m³"),
        ("NOx: Training-Mean Baseline", "0.0000", "11.04 mg/m³", "-0.6914", "13.76 mg/m³")
    ]
    for r_idx, row in enumerate(model_summary_data, start=1):
        for c_idx, val in enumerate(row):
            cell = t8.cell(r_idx, c_idx)
            cell.text = val
            cell.fill.solid()
            if r_idx in [1, 2]:
                cell.fill.fore_color.rgb = GREEN_CARD
            else:
                cell.fill.fore_color.rgb = WHITE if r_idx % 2 == 0 else LIGHT_BG
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(10)
            if r_idx in [1, 2]:
                p.font.bold = True
                p.font.color.rgb = DARK_BLUE
            else:
                p.font.color.rgb = TEXT_DARK

    # --- SLIDE 9: 6. Detailed Results - Energy Yield & PCR ---
    slide9 = prs.slides.add_slide(blank_layout)
    add_header(slide9, "6. Detailed Results: Energy Yield (TEY) & PCR Diagnostics", "PHASE 4: IN-DEPTH MODEL EVALUATION")
    card_d1 = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.4), Inches(5.6), Inches(5.6))
    card_d1.fill.solid()
    card_d1.fill.fore_color.rgb = LIGHT_BG
    card_d1.line.color.rgb = BORDER_COLOR
    tf_d1 = card_d1.text_frame
    tf_d1.word_wrap = True
    tf_d1.margin_left = Inches(0.25)
    tf_d1.margin_top = Inches(0.2)

    p = tf_d1.paragraphs[0]
    p.text = "TEY Model Diagnostics & Multicollinearity"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    tey_details = [
        ("Multicollinearity Diagnostics (Slide 115):", "• Extreme VIF scores diagnosed in operating telemetry: CDP (342.82), TIT (262.12), GTEP (260.47), TAT (133.21).\n• Predictor covariance eigen-decomposition reveals 4 PCs capture 90.76% of total feature variance."),
        ("OLS vs. PCR Tradeoff (Slides 70, 106):", "• Full OLS: Highest raw test accuracy (Test R² = 0.9501, RMSE = 3.35 MWh), but predictor weights suffer from severe collinear instability.\n• PCR (4 PCs): Achieves Test R² = 0.9397 (RMSE = 3.68 MWh) while projecting onto strictly orthogonal coordinates, guaranteeing stable physical interpretation."),
        ("Residual Analysis (Slides 104, 108):", "• Residuals $e_i = y_i - \\hat{y}_i$ are zero-centered with symmetric Gaussian distribution and homoscedastic variance across the full 100–180 MWh generator span.")
    ]
    for lbl, desc in tey_details:
        p_l = tf_d1.add_paragraph()
        p_l.text = f"\n{lbl}"
        p_l.font.bold = True
        p_l.font.size = Pt(11)
        p_l.font.color.rgb = ACCENT_BLUE
        p_d = tf_d1.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(10)
        p_d.font.color.rgb = TEXT_DARK
    if os.path.exists("figures/results_residuals_diagnostics.png"):
        slide9.shapes.add_picture("figures/results_residuals_diagnostics.png", Inches(6.7), Inches(1.4), width=Inches(5.8))

    # --- SLIDE 10: 6. Detailed Results - Emissions & Operational Regimes ---
    slide10 = prs.slides.add_slide(blank_layout)
    add_header(slide10, "6. Detailed Results: Emissions Modeling & Operational Regimes", "PHASE 4: IN-DEPTH MODEL EVALUATION")
    card_d2 = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.4), Inches(5.6), Inches(5.6))
    card_d2.fill.solid()
    card_d2.fill.fore_color.rgb = LIGHT_BG
    card_d2.line.color.rgb = BORDER_COLOR
    tf_d2 = card_d2.text_frame
    tf_d2.word_wrap = True
    tf_d2.margin_left = Inches(0.25)
    tf_d2.margin_top = Inches(0.2)

    p = tf_d2.paragraphs[0]
    p.text = "Emissions Modeling & K-Means Clusters"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    emiss_details = [
        ("CO Polynomial Enhancement (Slide 116):", "• Adding quadratic feature terms ($x \\to [x, x^2]$) increases Train R² to 0.6477 (RMSE = 1.36 mg/m³).\n• Captures the non-linear inflection below 1070°C where incomplete combustion surges exponentially."),
        ("NOx Domain Shift Analysis (*):", "• Negative out-of-sample Test R² (-1.0928) reflects a plant hardware/burner recalibration in 2014–15 (baseline NOx dropped from 68.8 to 58.5 mg/m³)."),
        ("Operational Regimes Profile (K=3, Silhouette = 0.3023):", "• Peak Load (4,309 hrs): TEY = 158.2 MWh, TIT = 1100°C, CDP = 13.8 mbar $\\implies$ CO = 1.00 mg/m³.\n• Nominal Load (7,820 hrs): TEY = 134.5 MWh, TIT = 1091°C, CDP = 12.2 mbar $\\implies$ CO = 1.48 mg/m³.\n• Part Load (10,062 hrs): TEY = 122.3 MWh, TIT = 1070°C, CDP = 11.2 mbar $\\implies$ CO = 3.30 mg/m³ (3.3x higher!).")
    ]
    for lbl, desc in emiss_details:
        p_l = tf_d2.add_paragraph()
        p_l.text = f"\n{lbl}"
        p_l.font.bold = True
        p_l.font.size = Pt(11)
        p_l.font.color.rgb = ACCENT_BLUE
        p_d = tf_d2.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(9.5)
        p_d.font.color.rgb = TEXT_DARK
    if os.path.exists("figures/results_kmeans_regimes_scatter.png"):
        slide10.shapes.add_picture("figures/results_kmeans_regimes_scatter.png", Inches(6.7), Inches(1.4), width=Inches(5.8))

    # --- SLIDE 11: 7. Conclusions & Industrial Insights ---
    slide11 = prs.slides.add_slide(blank_layout)
    add_header(slide11, "7. Conclusions & Engineering Recommendations", "PHASE 5: STRATEGIC INSIGHTS")
    concl_cards = [
        ("1. Methodological Justifications", 
         "• Multicollinearity: VIF > 250 proved that raw sensor telemetry is collinear. PCR on 4 orthogonal components captured 90.76% of variance, ensuring parameter stability.\n"
         "• Non-linear Physics: CO follows non-linear thermal oxidation kinetics; degree-2 polynomial expansion resolved the inflection point.\n"
         "• Unsupervised Regimes: K-Means (K=3) partitioned turbine operations into distinct thermodynamic states."),
        ("2. Operational Emission Control", 
         "• Avoid Low Firing Regimes: Part-load operation below 1070°C firing produces 3.3x higher CO emissions.\n"
         "• Peak Load Cleanliness: Operating at peak load (TIT ~1100°C) maximizes power yield (158.2 MWh) and minimizes CO (1.0 mg/m³).\n"
         "• Ambient Weather Compensation: Ambient temperature penalizes turbine yield (-0.58 correlation), justifying inlet air chilling systems during summer."),
        ("3. Digital Twin & Deployment", 
         "• Real-Time Monitoring: The lightweight 4-component PCR model can run in real-time on edge turbine PLCs for anomaly detection.\n"
         "• Sensor Redundancy: High collinearity allows synthetic sensor estimation if a pressure transducer fails in operation.\n"
         "• Regulatory Compliance: Predictive emission models enable preemptive flue gas compliance adjustments.")
    ]
    for i, (title, desc) in enumerate(concl_cards):
        left = Inches(0.8 + i * 4.0)
        card_c = slide11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.5), Inches(3.7), Inches(5.5))
        card_c.fill.solid()
        card_c.fill.fore_color.rgb = LIGHT_BG
        card_c.line.color.rgb = BORDER_COLOR
        tf_c = card_c.text_frame
        tf_c.word_wrap = True
        tf_c.margin_left = Inches(0.25)
        tf_c.margin_top = Inches(0.2)
        p = tf_c.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE
        p_desc = tf_c.add_paragraph()
        p_desc.text = f"\n{desc}"
        p_desc.font.size = Pt(10.5)
        p_desc.font.color.rgb = TEXT_DARK

    # --- SLIDE 12: Project Summary & Q&A ---
    slide12 = prs.slides.add_slide(blank_layout)
    bg12 = slide12.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg12.fill.solid()
    bg12.fill.fore_color.rgb = DARK_BLUE
    bg12.line.color.rgb = DARK_BLUE

    tbox12 = slide12.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.333), Inches(4.0))
    tf12 = tbox12.text_frame
    tf12.word_wrap = True
    p = tf12.paragraphs[0]
    p.text = "Thank You!"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE

    p2 = tf12.add_paragraph()
    p2.text = "Questions & Technical Discussion\n"
    p2.font.size = Pt(22)
    p2.font.color.rgb = RGBColor(180, 215, 235)

    p3 = tf12.add_paragraph()
    p3.text = (
        "• Supplementary Material: Full Python Jupyter Notebook (gas_turbine_analysis.ipynb)\n"
        "• Dataset: 36,733 hourly records across 2011–2015 (Turkey Gas Turbine)\n"
        "• All methods strictly conform to the Data Analytics course syllabus."
    )
    p3.font.size = Pt(13)
    p3.font.color.rgb = RGBColor(210, 225, 240)

    prs.save("presentation.pptx")
    print("presentation.pptx updated cleanly with all 12 slides.")

update_presentation()
